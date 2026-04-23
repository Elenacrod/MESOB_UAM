"""
CRIM Intervals - Presentation Generator
18-slide PowerPoint in Spanish for musicology students
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import math

OUTPUT_PATH = "C:/Users/MC.5055521/Desktop/Claude code/CRIM_Intervals_Musicologia.pptx"

# ── Color palette ──────────────────────────────────────────────────────────────
BURGUNDY      = RGBColor(0x6D, 0x2E, 0x46)
CREAM         = RGBColor(0xEC, 0xE2, 0xD0)
DUSTY_ROSE    = RGBColor(0xA2, 0x67, 0x69)
DARK_BURGUNDY = RGBColor(0x4A, 0x1E, 0x30)
LIGHT_CREAM   = RGBColor(0xF5, 0xEF, 0xE6)
MED_BURGUNDY  = RGBColor(0x8B, 0x3A, 0x54)
DARK_TEXT     = RGBColor(0x2C, 0x10, 0x18)
ROSE_LIGHT    = RGBColor(0xC8, 0x97, 0x9A)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)

# ── Setup ──────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)

BLANK_LAYOUT = prs.slide_layouts[6]  # blank

def in_(x): return Inches(x)
def pt_(x): return Pt(x)

# ── Low-level helpers ──────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=0):
    from pptx.util import Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        in_(x), in_(y), in_(w), in_(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color and line_width > 0:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_oval(slide, x, y, w, h, fill_color, line_color=None, line_width=0, transparency=0):
    shape = slide.shapes.add_shape(
        9,  # MSO_AUTO_SHAPE_TYPE.OVAL = 9
        in_(x), in_(y), in_(w), in_(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if transparency > 0:
        shape.fill.fore_color.theme_color  # ensure theme access
        from pptx.oxml.ns import qn
        from lxml import etree
        # Set alpha via XML
        sp_pr = shape._element
        solid_fill = sp_pr.find('.//' + qn('a:solidFill'))
        if solid_fill is not None:
            srgb = solid_fill.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha = srgb.find(qn('a:alpha'))
                if alpha is None:
                    alpha = etree.SubElement(srgb, qn('a:alpha'))
                alpha.set('val', str(int((100 - transparency) * 1000)))
    if line_color and line_width > 0:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def set_bg(slide, color):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, x, y, w, h, text, font_size, font_name, color,
                bold=False, italic=False, align=PP_ALIGN.LEFT, valign=None,
                wrap=True, margin_l=None, margin_r=None, margin_t=None, margin_b=None):
    txBox = slide.shapes.add_textbox(in_(x), in_(y), in_(w), in_(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    # Set margins
    if margin_l is not None: tf.margin_left   = in_(margin_l)
    if margin_r is not None: tf.margin_right  = in_(margin_r)
    if margin_t is not None: tf.margin_top    = in_(margin_t)
    if margin_b is not None: tf.margin_bottom = in_(margin_b)

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = pt_(font_size)
    run.font.name = font_name
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return txBox

def add_textbox_multiline(slide, x, y, w, h, lines, font_size, font_name, color,
                          bold=False, italic=False, align=PP_ALIGN.LEFT,
                          margin_l=0.05, margin_t=0.05):
    """lines: list of (text, bold, italic, bullet) tuples"""
    txBox = slide.shapes.add_textbox(in_(x), in_(y), in_(w), in_(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left   = in_(margin_l)
    tf.margin_top    = in_(margin_t)
    tf.margin_right  = in_(0.05)
    tf.margin_bottom = in_(0.05)

    first = True
    for (text, b, it, bullet) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if bullet:
            from pptx.oxml.ns import qn
            from lxml import etree
            pPr = p._pPr
            if pPr is None:
                pPr = p._p.get_or_add_pPr()
            buChar = etree.SubElement(pPr, qn('a:buChar'))
            buChar.set('char', '•')
            buFont = etree.SubElement(pPr, qn('a:buFont'))
            buFont.set('typeface', 'Arial')
            indent = etree.SubElement(pPr, qn('a:indent')) if pPr.find(qn('a:indent')) is None else pPr.find(qn('a:indent'))
            indent.set('val', '-342900')  # ~0.36 cm hanging indent
            marL = pPr.get('marL')
            pPr.set('marL', '342900')
        run = p.add_run()
        run.text = text
        run.font.size = pt_(font_size)
        run.font.name = font_name
        run.font.color.rgb = color
        run.font.bold = b if b is not None else bold
        run.font.italic = it if it is not None else italic
    return txBox

def add_header_band(slide, title, color=None, font_size=30):
    """Add a full-width colored header band with title."""
    if color is None:
        color = BURGUNDY
    add_rect(slide, 0, 0, 10, 1.0, color)
    add_textbox(slide, 0.3, 0.08, 9.4, 0.85, title,
                font_size, "Georgia", CREAM, bold=True,
                align=PP_ALIGN.LEFT, margin_l=0.15, margin_t=0.05)

def add_accent_bar(slide):
    """Left burgundy accent bar."""
    add_rect(slide, 0, 0, 0.12, 5.625, BURGUNDY)

def add_corner_ornament(slide):
    """Decorative overlapping circles for dark slides."""
    add_oval(slide, 7.8, -1.0, 3.2, 3.2, BURGUNDY, transparency=60)
    add_oval(slide, 8.6, -0.3, 2.0, 2.0, DUSTY_ROSE, transparency=70)

def add_badge(slide, text, x, y, w=3.4):
    add_rect(slide, x, y, w, 0.38, BURGUNDY)
    add_textbox(slide, x, y, w, 0.38, text, 12, "Calibri", CREAM,
                bold=True, align=PP_ALIGN.CENTER, margin_l=0, margin_t=0.05)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide1, DARK_BURGUNDY)

# Decorative circles
add_oval(slide1, 6.5, 0.2, 4.2, 4.2, BURGUNDY, transparency=30)
add_oval(slide1, 7.2, 1.0, 2.8, 2.8, DUSTY_ROSE, transparency=50)

# Staff lines
for i in range(5):
    add_rect(slide1, 6.2, 1.5 + i * 0.28, 3.5, 0.04, ROSE_LIGHT)

# Note head + stem
add_oval(slide1, 7.6, 2.1, 0.55, 0.42, CREAM)
add_rect(slide1, 8.12, 1.05, 0.05, 1.1, CREAM)

# Left accent
add_rect(slide1, 0, 0, 0.18, 5.625, DUSTY_ROSE)

add_textbox(slide1, 0.5, 0.7, 6.0, 1.1, "CRIM Intervals",
            48, "Georgia", CREAM, bold=True)

add_textbox(slide1, 0.5, 1.85, 6.0, 0.8,
            "Herramientas de Análisis para la Polifonía Renacentista",
            19, "Georgia", ROSE_LIGHT, italic=True)

add_textbox(slide1, 0.5, 2.72, 6.0, 0.55,
            "Una introducción para estudiantes de musicología",
            15, "Calibri", CREAM)

add_rect(slide1, 0.5, 3.42, 5.5, 0.04, DUSTY_ROSE)

add_textbox(slide1, 0.5, 5.05, 9.0, 0.38,
            "Citations: The Renaissance Imitation Mass  ·  crimproject.org",
            11, "Calibri", ROSE_LIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — ¿Qué es el Proyecto CRIM?
# ══════════════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide2, CREAM)
add_accent_bar(slide2)
add_header_band(slide2, "¿Qué es el Proyecto CRIM?")

add_textbox(slide2, 0.3, 1.15, 4.8, 0.42, "Sobre el Proyecto",
            15, "Georgia", BURGUNDY, bold=True)

add_textbox_multiline(slide2, 0.3, 1.62, 4.7, 2.1,
    [
        ("CRIM = Citations: The Renaissance Imitation Mass", True, False, True),
        ("Estudia cómo compositores del s. XVI transformaron obras en ciclos de Misa", False, False, True),
        ("Dirigido por el Prof. Richard Freedman (Haverford College)", False, False, True),
    ],
    14, "Calibri", DARK_TEXT)

# Stat boxes
stats = [
    ("50+",    "investigadores internacionales", BURGUNDY),
    ("~3.000", "relaciones documentadas",        DUSTY_ROSE),
    ("2",      "instituciones principales",      BURGUNDY),
]
for i, (num, label, col) in enumerate(stats):
    bx, by = 5.4, 1.18 + i * 1.15
    add_rect(slide2, bx, by, 4.2, 0.95, col)
    num_font = 22 if len(num) > 3 else 26
    add_textbox(slide2, bx + 0.1, by, 1.7, 0.95, num,
                num_font, "Georgia", CREAM, bold=True, margin_l=0.05, margin_t=0.22)
    add_textbox(slide2, bx + 1.82, by, 2.28, 0.95, label,
                12, "Calibri", CREAM, margin_l=0, margin_t=0.3)

add_textbox(slide2, 0.3, 4.62, 9.4, 0.5,
            "Haverford College  +  Centre d'Études Supérieures de la Renaissance (Tours, Francia)",
            11, "Calibri", DUSTY_ROSE, italic=True)

add_oval(slide2, 8.5, 3.85, 1.7, 1.7, BURGUNDY, line_color=DUSTY_ROSE, line_width=1, transparency=82)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Aplicación Streamlit
# ══════════════════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide3, CREAM)
add_accent_bar(slide3)
add_header_band(slide3, "Acceso Sin Código: La Aplicación Streamlit")

rows3 = [
    (BURGUNDY,   "1", "Sin programación",        "Interfaz web intuitiva, sin necesidad de conocimientos técnicos"),
    (DUSTY_ROSE, "2", "Formatos abiertos",        "Compatible con MEI, MusicXML y MIDI"),
    (BURGUNDY,   "3", "Resultados descargables",  "Tablas, gráficos y CSV para publicaciones"),
    (DUSTY_ROSE, "4", "Corpus integrado",          "Acceso directo al corpus CRIM + tus propias partituras"),
]

for i, (col, num, title, desc) in enumerate(rows3):
    ry = 1.15 + i * 0.9
    add_oval(slide3, 0.3, ry + 0.05, 0.62, 0.62, col)
    add_textbox(slide3, 0.3, ry + 0.05, 0.62, 0.62, num,
                16, "Georgia", CREAM, bold=True, align=PP_ALIGN.CENTER,
                margin_l=0, margin_t=0.12)
    add_textbox(slide3, 1.1, ry + 0.02, 3.5, 0.35, title,
                14, "Georgia", BURGUNDY, bold=True, margin_t=0.02)
    add_textbox(slide3, 1.1, ry + 0.38, 8.6, 0.35, desc,
                13, "Calibri", DARK_TEXT, margin_t=0.02)

add_rect(slide3, 0.3, 4.88, 9.4, 0.44, DARK_BURGUNDY)
add_textbox(slide3, 0.3, 4.88, 9.4, 0.44,
            "Built on:  Python  ·  CRIM Intervals Library  ·  music21  ·  Streamlit",
            12, "Calibri", CREAM, bold=True, align=PP_ALIGN.CENTER,
            margin_l=0, margin_t=0.07)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — El Repertorio (dark)
# ══════════════════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide4, DARK_BURGUNDY)
add_corner_ornament(slide4)
add_rect(slide4, 0, 0, 0.25, 5.625, DUSTY_ROSE)

add_textbox(slide4, 0.5, 0.4, 7.5, 0.7, "El Repertorio",
            40, "Georgia", CREAM, bold=True)
add_textbox(slide4, 0.5, 1.1, 7.5, 0.45, "La Misa Imitación Renacentista",
            18, "Georgia", ROSE_LIGHT, italic=True)

bullets4 = [
    "Género clave del siglo XVI: la Misa parodia o de imitación",
    "Técnica compositiva: transformar motetes, chansons y otras piezas en ciclos de Misa de 5 movimientos",
    "Compositores: Josquin · Palestrina · Lassus · Morales · Victoria",
    "Técnicas contrapuntísticas: fugas, dúos imitativos, entradas periódicas",
    "El análisis computacional permite estudiar este repertorio a escala masiva",
]
for i, b in enumerate(bullets4):
    by = 1.72 + i * 0.63
    add_rect(slide4, 0.5, by + 0.13, 0.22, 0.22, DUSTY_ROSE)
    add_textbox(slide4, 0.9, by, 8.8, 0.52, b, 14, "Calibri", CREAM, margin_t=0.08)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Herramienta 1: Intervalos Melódicos
# ══════════════════════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide5, CREAM)
add_accent_bar(slide5)
add_header_band(slide5, "Herramienta 1: Análisis de Intervalos Melódicos", font_size=26)
add_badge(slide5, "piece.melodic()", 0.3, 1.1)

add_textbox(slide5, 0.3, 1.62, 4.5, 0.4, "¿Qué hace?",
            15, "Georgia", BURGUNDY, bold=True)
add_textbox_multiline(slide5, 0.3, 2.08, 4.5, 1.8, [
    ("Identifica y cuantifica todos los intervalos melódicos de cada voz", False, False, True),
    ("Opciones: diatónico o cromático", False, False, True),
    ("Con o sin calidad interválica (mayor, menor, perfecto)", False, False, True),
], 14, "Calibri", DARK_TEXT)

# Vertical divider
add_rect(slide5, 5.1, 1.1, 0.06, 3.8, DUSTY_ROSE)

add_textbox(slide5, 5.3, 1.62, 4.4, 0.4, "Utilidad para la polifonía renacentista",
            15, "Georgia", BURGUNDY, bold=True)
add_textbox_multiline(slide5, 5.3, 2.08, 4.4, 2.2, [
    ("Detectar patrones motívicos recurrentes", False, False, True),
    ("Comparar los perfiles melódicos entre voces", False, False, True),
    ("Identificar material prestado (citación melódica)", False, False, True),
    ("Estudiar el estilo de cada compositor", False, False, True),
], 14, "Calibri", DARK_TEXT)

add_oval(slide5, 8.6, 4.1, 1.2, 1.2, BURGUNDY, line_color=DUSTY_ROSE, line_width=1.5, transparency=78)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Herramienta 2: Intervalos Armónicos
# ══════════════════════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide6, CREAM)
add_accent_bar(slide6)
add_header_band(slide6, "Herramienta 2: Análisis de Intervalos Armónicos", DUSTY_ROSE, font_size=26)
add_badge(slide6, "piece.harmonic()", 0.3, 1.1)

add_textbox(slide6, 0.3, 1.62, 4.5, 0.4, "¿Qué hace?",
            15, "Georgia", BURGUNDY, bold=True)
add_textbox_multiline(slide6, 0.3, 2.08, 4.5, 1.5, [
    ("Examina los intervalos verticales entre cada par de voces", False, False, True),
    ("Mide la distancia interválica en cada punto de la textura", False, False, True),
], 14, "Calibri", DARK_TEXT)

# Mini bar chart visualization
bar_data = [("8va", 3.2, BURGUNDY), ("5ta", 2.6, DUSTY_ROSE), ("3ra", 2.0, MED_BURGUNDY), ("6ta", 2.3, ROSE_LIGHT)]
for i, (label, bw, col) in enumerate(bar_data):
    by = 3.55 + i * 0.36
    add_rect(slide6, 0.3, by, bw, 0.28, col)
    add_textbox(slide6, 0.3 + bw + 0.1, by, 0.8, 0.28, label,
                11, "Calibri", DARK_TEXT, margin_t=0)

add_rect(slide6, 5.1, 1.1, 0.06, 3.8, DUSTY_ROSE)

add_textbox(slide6, 5.3, 1.62, 4.4, 0.4, "Utilidad para la polifonía renacentista",
            15, "Georgia", BURGUNDY, bold=True)
add_textbox_multiline(slide6, 5.3, 2.08, 4.4, 2.2, [
    ("Detectar consonancias y disonancias características del estilo", False, False, True),
    ("Analizar el tratamiento del contrapunto estricto", False, False, True),
    ("Identificar progresiones armónicas características", False, False, True),
    ("Estudiar la densidad contrapuntística de la textura", False, False, True),
], 14, "Calibri", DARK_TEXT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Herramienta 3: N-gramas
# ══════════════════════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide7, CREAM)
add_accent_bar(slide7)
# Custom taller header for this long title
add_rect(slide7, 0, 0, 10, 1.2, BURGUNDY)
add_textbox(slide7, 0.3, 0.08, 9.4, 1.05,
            "Herramienta 3: N-gramas y Módulos Contrapuntísticos",
            24, "Georgia", CREAM, bold=True,
            align=PP_ALIGN.LEFT, margin_l=0.15, margin_t=0.22)

add_rect(slide7, 0.3, 1.3, 9.4, 0.65, DUSTY_ROSE)
add_textbox(slide7, 0.3, 1.3, 9.4, 0.65, "El corazón analítico de CRIM Intervals",
            18, "Georgia", CREAM, bold=True, italic=True, align=PP_ALIGN.CENTER,
            margin_l=0, margin_t=0.12)

add_textbox(slide7, 0.3, 2.08, 5.0, 0.38, "Un módulo contrapuntístico combina:",
            14, "Calibri", DARK_TEXT, bold=True)

add_rect(slide7, 0.3, 2.32, 3.8, 0.72, BURGUNDY)
add_textbox(slide7, 0.3, 2.32, 3.8, 0.72, "Intervalos verticales entre voces",
            13, "Calibri", CREAM, bold=True, align=PP_ALIGN.CENTER,
            margin_l=0, margin_t=0.2)

add_textbox(slide7, 4.05, 2.32, 0.3, 0.72, "+", 22, "Georgia", BURGUNDY,
            bold=True, align=PP_ALIGN.CENTER, margin_l=0, margin_t=0.18)

add_rect(slide7, 4.3, 2.32, 3.8, 0.72, DUSTY_ROSE)
add_textbox(slide7, 4.3, 2.32, 3.8, 0.72, "Movimiento melódico de la voz inferior",
            13, "Calibri", CREAM, bold=True, align=PP_ALIGN.CENTER,
            margin_l=0, margin_t=0.2)

add_rect(slide7, 0.3, 3.18, 5.2, 0.65, DARK_BURGUNDY)
add_textbox(slide7, 0.3, 3.18, 5.2, 0.65, '"7_Held · 6_-2 · 8"',
            18, "Calibri", CREAM, bold=True, align=PP_ALIGN.CENTER,
            margin_l=0, margin_t=0.14)

add_textbox(slide7, 0.3, 3.95, 9.4, 0.42,
            "Intervalos verticales 7→6→8 + nota ligada + descenso de 2a = fórmula cadencial típica",
            13, "Calibri", DARK_TEXT, italic=True)

add_rect(slide7, 0.3, 4.5, 9.4, 0.75, LIGHT_CREAM)
# Border via line
shape7 = slide7.shapes[-1]
shape7.line.color.rgb = DUSTY_ROSE
shape7.line.width = Pt(1)
add_textbox(slide7, 0.3, 4.5, 9.4, 0.75,
            "Utilidad: Detectar citas y transformaciones  ·  Comparar técnicas entre compositores",
            13, "Calibri", BURGUNDY, bold=True, align=PP_ALIGN.CENTER,
            margin_l=0, margin_t=0.2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Herramienta 4: Cadencias
# ══════════════════════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide8, CREAM)
add_accent_bar(slide8)
add_header_band(slide8, "Herramienta 4: Detección de Cadencias", DUSTY_ROSE)
add_badge(slide8, "piece.cadences()", 0.3, 1.1)

cols8 = [
    ("Tipos detectados",     ["Cadencia auténtica", "Cadencia frigia", "Cadencia evitada", "Cadencia de engaño"], 0.25, BURGUNDY),
    ("Información extraída", ["Voz que realiza la cláusula", "Tono de la cadencia", "Posición métrica y formal"], 3.6, DUSTY_ROSE),
    ("Utilidad musicológica",["Trazar el perfil tonal (modal)", "Comparar modelo vs. Misa", "Analizar la articulación formal"], 6.7, MED_BURGUNDY),
]
for title, items, cx, col in cols8:
    add_rect(slide8, cx, 1.65, 3.0, 0.45, col)
    add_textbox(slide8, cx, 1.65, 3.0, 0.45, title, 13, "Georgia", CREAM,
                bold=True, align=PP_ALIGN.CENTER, margin_l=0, margin_t=0.1)
    for j, item in enumerate(items):
        iy = 2.18 + j * 0.75
        add_rect(slide8, cx, iy, 3.0, 0.65, LIGHT_CREAM)
        shape8 = slide8.shapes[-1]
        shape8.line.color.rgb = col
        shape8.line.width = Pt(1)
        add_textbox(slide8, cx + 0.08, iy, 2.84, 0.65, item, 13, "Calibri", DARK_TEXT,
                    align=PP_ALIGN.CENTER, margin_l=0, margin_t=0.18)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Herramienta 5: Tipos de Presentación
# ══════════════════════════════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide9, CREAM)
add_accent_bar(slide9)
add_header_band(slide9, "Herramienta 5: Tipos de Presentación Imitativa", font_size=26)
add_badge(slide9, "piece.presentationTypes()", 0.3, 1.1)

types9 = [
    ("FUGA",                "Canon o quasi-canon entre voces"),
    ("DÚO IMITATIVO",       "Dos voces en imitación estricta"),
    ("DÚO NO IMITATIVO",    "Dos voces en contrapunto libre"),
    ("ENTRADAS PERIÓDICAS", "Todas las voces con el mismo motivo"),
    ("HOMOFONÍA RÍTMICA",   "piece.homorhythm() — Movimiento sincrónico de voces"),
]
for i, (name, desc) in enumerate(types9):
    ty = 1.62 + i * 0.63
    col = BURGUNDY if i % 2 == 0 else DUSTY_ROSE
    add_rect(slide9, 0.3, ty + 0.08, 0.44, 0.42, col)
    add_textbox(slide9, 0.9, ty, 2.8, 0.55, name, 14, "Georgia", BURGUNDY,
                bold=True, margin_t=0.1)
    add_textbox(slide9, 3.7, ty, 5.9, 0.55, "— " + desc, 13, "Calibri", DARK_TEXT,
                margin_t=0.1)

add_rect(slide9, 0.3, 4.9, 9.4, 0.44, DARK_BURGUNDY)
add_textbox(slide9, 0.3, 4.9, 9.4, 0.44,
            "Utilidad: Mapear la arquitectura contrapuntística  ·  Comparar uso del material en la Misa",
            12, "Calibri", CREAM, bold=True, align=PP_ALIGN.CENTER,
            margin_l=0, margin_t=0.08)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Herramienta 6: Heatmaps
# ══════════════════════════════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide10, CREAM)
add_accent_bar(slide10)
add_rect(slide10, 0, 0, 10, 1.0, DUSTY_ROSE)
add_textbox(slide10, 0.3, 0.08, 9.4, 0.88,
            "Herramienta 6: Mapas de Calor (Heatmaps)",
            26, "Georgia", CREAM, bold=True,
            align=PP_ALIGN.LEFT, margin_l=0.15, margin_t=0.16)

heat_colors = [LIGHT_CREAM, ROSE_LIGHT, DUSTY_ROSE, MED_BURGUNDY, BURGUNDY, DARK_BURGUNDY]
heatX, heatY = 0.2, 1.15
cellW, cellH = 0.41, 0.32
heatRows, heatCols = 7, 12
for r in range(heatRows):
    for c in range(heatCols):
        val = math.sin(c * 0.5 + r * 0.7) * 0.5 + math.cos(c * 0.3 - r * 0.4) * 0.4
        intensity = max(0, min(5, int((val + 1) * 2.5)))
        add_rect(slide10,
                 heatX + c * cellW, heatY + r * cellH,
                 cellW - 0.03, cellH - 0.03,
                 heat_colors[intensity])

add_textbox(slide10, 5.5, 1.2, 4.2, 0.42, "¿Qué muestran?",
            15, "Georgia", BURGUNDY, bold=True)
add_textbox_multiline(slide10, 5.5, 1.67, 4.2, 1.4, [
    ("Eje horizontal: posición en la partitura (compases)", False, False, True),
    ("Eje vertical: tipo de patrón o voz", False, False, True),
    ("Color: frecuencia o densidad del patrón", False, False, True),
], 13, "Calibri", DARK_TEXT)

add_textbox(slide10, 5.5, 3.2, 4.2, 0.42, "Utilidad",
            15, "Georgia", BURGUNDY, bold=True)
add_textbox_multiline(slide10, 5.5, 3.67, 4.2, 1.7, [
    ("Ver dónde se concentran las citas y transformaciones", False, False, True),
    ("Comparar arquitectura formal de modelo y Misa", False, False, True),
    ("Identificar secciones de imitación densa vs. textura libre", False, False, True),
], 13, "Calibri", DARK_TEXT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Herramienta 7: Radar Plots
# ══════════════════════════════════════════════════════════════════════════════
slide11 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide11, CREAM)
add_accent_bar(slide11)
add_header_band(slide11, "Herramienta 7: Radar Plots y Gráficos de Cadencias", font_size=24)

add_textbox(slide11, 0.3, 1.15, 4.5, 0.42, "Radar Plot (diagrama radial)",
            15, "Georgia", BURGUNDY, bold=True)
add_textbox_multiline(slide11, 0.3, 1.62, 4.5, 1.4, [
    ("Muestra el perfil tonal de la obra", False, False, True),
    ('Cada "radio" = un grado o tono cadencial', False, False, True),
    ("Compara el perfil modal de varias obras", False, False, True),
], 13, "Calibri", DARK_TEXT)

# Decorative radar
cx, cy = 2.5, 4.0
ring_colors = [ROSE_LIGHT, DUSTY_ROSE, BURGUNDY]
for ring_i, rcolor in enumerate(ring_colors):
    rr = (ring_i + 1) * 0.38
    add_oval(slide11, cx - rr, cy - rr, rr * 2, rr * 2, rcolor,
             line_color=BURGUNDY, line_width=1, transparency=65 + ring_i * 4)

add_rect(slide11, 5.1, 1.1, 0.06, 4.1, DUSTY_ROSE)

add_textbox(slide11, 5.3, 1.15, 4.4, 0.42, "Progress Charts (Cadencias)",
            15, "Georgia", BURGUNDY, bold=True)
add_textbox_multiline(slide11, 5.3, 1.62, 4.4, 1.4, [
    ("Cadencias en secuencia temporal", False, False, True),
    ("Visualiza la planificación tonal", False, False, True),
    ("Analiza la articulación formal de los movimientos", False, False, True),
], 13, "Calibri", DARK_TEXT)

# Mini bar chart
bar11 = [0.3, 0.7, 0.5, 0.9, 0.4, 0.6, 0.8, 0.35, 0.55, 0.75]
bar_cols = [BURGUNDY, DUSTY_ROSE] * 5
for i, (bh, bcol) in enumerate(zip(bar11, bar_cols)):
    add_rect(slide11, 5.3 + i * 0.41, 4.5 - bh * 1.2, 0.34, bh * 1.2, bcol)
add_rect(slide11, 5.3, 4.5, 4.1, 0.04, DARK_TEXT)
add_textbox(slide11, 5.3, 4.56, 4.1, 0.3, "Secuencia temporal de cadencias",
            10, "Calibri", DUSTY_ROSE, italic=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Herramienta 8: Network Diagrams
# ══════════════════════════════════════════════════════════════════════════════
slide12 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide12, CREAM)
add_accent_bar(slide12)
add_header_band(slide12, "Herramienta 8: Diagramas de Red", DUSTY_ROSE)

nodes12 = [
    (0.4,  2.8, 0.3, BURGUNDY,  "Motete"),
    (2.0,  1.6, 0.25, DUSTY_ROSE, "Kyrie"),
    (2.0,  4.0, 0.25, DUSTY_ROSE, "Gloria"),
    (3.6,  1.1, 0.2, MED_BURGUNDY,"Credo"),
    (3.6,  2.8, 0.22, MED_BURGUNDY,"Sanctus"),
    (3.6,  4.5, 0.2, DUSTY_ROSE,"Agnus"),
]

# Draw edges first
edges12 = [(0,1),(0,2),(0,4),(1,3),(1,4),(2,4),(2,5),(3,4),(4,5)]
for a, b in edges12:
    na, nb = nodes12[a], nodes12[b]
    # Draw as a thin rectangle approximating a line
    ax, ay, ar = na[0] + na[2], na[1] + na[2], na[2]
    bx, by2, br = nb[0] + nb[2], nb[1] + nb[2], nb[2]
    # Use a line shape via connector - approximate with very thin rect
    from pptx.util import Emu
    from pptx.oxml.ns import qn
    from lxml import etree
    # Use add_connector from shapes
    cnx = slide12.shapes.add_connector(1, in_(ax), in_(ay), in_(bx), in_(by2))
    cnx.line.color.rgb = DUSTY_ROSE
    cnx.line.width = Pt(1.2)

# Draw nodes
for nx, ny, nr, ncol, nlabel in nodes12:
    add_oval(slide12, nx, ny, nr * 2, nr * 2, ncol, line_color=CREAM, line_width=1.5)
    add_textbox(slide12, nx - 0.1, ny + nr * 2 + 0.04, 0.85, 0.28,
                nlabel, 10, "Calibri", DARK_TEXT, align=PP_ALIGN.CENTER, margin_t=0)

add_textbox(slide12, 5.0, 1.2, 4.7, 0.8,
            "Nodos = obras o motivos musicales\nAristas = similitudes, citas o relaciones",
            14, "Calibri", DARK_TEXT)

add_textbox(slide12, 5.0, 2.1, 4.7, 0.42, "Utilidad", 15, "Georgia", BURGUNDY, bold=True)
add_textbox_multiline(slide12, 5.0, 2.58, 4.7, 2.4, [
    ("Visualizar la red de influencias entre modelo y misas", False, False, True),
    ("Detectar agrupaciones de obras por patrones compartidos", False, False, True),
    ("Mostrar la genealogía de un motivo a través del corpus", False, False, True),
    ("Investigar la transmisión musical y el estilo compositivo", False, False, True),
], 13, "Calibri", DARK_TEXT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Corpus Tools (dark)
# ══════════════════════════════════════════════════════════════════════════════
slide13 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide13, DARK_BURGUNDY)
add_corner_ornament(slide13)
add_rect(slide13, 0, 0, 0.25, 5.625, DUSTY_ROSE)

add_textbox(slide13, 0.5, 0.3, 9.0, 0.65, "Herramienta 9: Análisis de Corpus",
            28, "Georgia", CREAM, bold=True)
add_textbox(slide13, 0.5, 1.02, 8.5, 0.42, "Del análisis individual al estudio a gran escala",
            14, "Georgia", ROSE_LIGHT, italic=True)

bullets13 = [
    "Aplica todos los métodos anteriores sobre múltiples obras simultáneamente",
    "Resultados agregados: distribución de intervalos, perfiles tonales, frecuencia de cadencias",
    "Comparación estadística entre compositores, géneros y períodos",
]
for i, b in enumerate(bullets13):
    by = 1.65 + i * 0.63
    add_rect(slide13, 0.5, by + 0.12, 0.22, 0.22, DUSTY_ROSE)
    add_textbox(slide13, 0.9, by, 8.8, 0.52, b, 14, "Calibri", CREAM, margin_t=0.08)

callouts13 = [
    ('"Close Reading"',   "Análisis detallado de una obra"),
    ('"Distant Reading"', "Patrones en cientos de obras a la vez"),
]
for i, (t, s) in enumerate(callouts13):
    bx = 0.5 + i * 4.8
    col = BURGUNDY if i == 0 else MED_BURGUNDY
    add_rect(slide13, bx, 3.92, 4.3, 1.1, col)
    shape_c = slide13.shapes[-1]
    shape_c.line.color.rgb = DUSTY_ROSE
    shape_c.line.width = Pt(1)
    add_textbox(slide13, bx + 0.1, 3.97, 4.1, 0.5, t,
                17, "Georgia", CREAM, bold=True, align=PP_ALIGN.CENTER,
                margin_l=0, margin_t=0.06)
    add_textbox(slide13, bx + 0.1, 4.47, 4.1, 0.45, s,
                12, "Calibri", ROSE_LIGHT, align=PP_ALIGN.CENTER,
                margin_l=0, margin_t=0.06)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Formatos de Archivo
# ══════════════════════════════════════════════════════════════════════════════
slide14 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide14, CREAM)
add_accent_bar(slide14)
add_header_band(slide14, "Formatos de Archivo Compatibles")

cards14 = [
    (BURGUNDY,   "MEI",      ["Music Encoding Initiative", "Estándar académico internacional", "Usado por el corpus CRIM"],          0.25, True,  CREAM,  CREAM),
    (DUSTY_ROSE, "MusicXML", ["Formato universal de notación", "Compatible con Sibelius, Finale,", "MuseScore, Dorico"],              3.6,  True,  CREAM,  CREAM),
    (LIGHT_CREAM,"MIDI",     ["Formato básico de datos MIDI", "Para análisis exploratorio inicial"],                                   6.88, False, DARK_TEXT, BURGUNDY),
]
for card_color, card_title, card_lines, cx, is_dark, tc, ttc in cards14:
    add_rect(slide14, cx, 1.2, 3.1, 3.05, card_color)
    if not is_dark:
        shape14 = slide14.shapes[-1]
        shape14.line.color.rgb = DUSTY_ROSE
        shape14.line.width = Pt(2)

    icon_bg = DARK_BURGUNDY if is_dark else DUSTY_ROSE
    add_oval(slide14, cx + 1.05, 1.35, 1.0, 1.0, icon_bg)

    letters = {"MEI": "M", "MusicXML": "X", "MIDI": "~"}
    add_textbox(slide14, cx + 1.05, 1.35, 1.0, 1.0, letters[card_title],
                22, "Georgia", CREAM, bold=True, align=PP_ALIGN.CENTER,
                margin_l=0, margin_t=0.2)
    add_textbox(slide14, cx + 0.1, 2.45, 2.9, 0.5, card_title,
                22, "Georgia", ttc, bold=True, align=PP_ALIGN.CENTER,
                margin_l=0, margin_t=0.04)
    for li, line in enumerate(card_lines):
        add_textbox(slide14, cx + 0.1, 3.0 + li * 0.32, 2.9, 0.3,
                    line, 12, "Calibri", tc, align=PP_ALIGN.CENTER,
                    margin_l=0, margin_t=0)

add_rect(slide14, 0.25, 4.55, 9.5, 0.72, LIGHT_CREAM)
shape14b = slide14.shapes[-1]
shape14b.line.color.rgb = DUSTY_ROSE
shape14b.line.width = Pt(1)
add_textbox(slide14, 0.25, 4.55, 9.5, 0.72,
            "El corpus CRIM ofrece ediciones de alta calidad en formato MEI  ·  También puedes importar tus propias partituras",
            12, "Calibri", BURGUNDY, italic=True, align=PP_ALIGN.CENTER,
            margin_l=0, margin_t=0.18)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Flujo de Trabajo
# ══════════════════════════════════════════════════════════════════════════════
slide15 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide15, CREAM)
add_accent_bar(slide15)
add_header_band(slide15, "Flujo de Trabajo para Estudiantes")

steps15 = [
    ("1", "Acceder",      "intervals-streamlit.crimproject.org"),
    ("2", "Cargar",       "Una obra del corpus CRIM o tu propio archivo MEI/MusicXML"),
    ("3", "Seleccionar",  "La herramienta de análisis deseada"),
    ("4", "Configurar",   "Parámetros: tipo de intervalo, voces, n-grama, etc."),
    ("5", "Visualizar",   "Resultados en tablas y gráficos interactivos"),
    ("6", "Descargar",    "Datos en CSV para publicaciones académicas"),
]
for i, (num, title, desc) in enumerate(steps15):
    col = i % 3
    row = i // 3
    sx = 0.3 + col * 3.22
    sy = 1.15 + row * 2.1
    step_col = BURGUNDY if col % 2 == 0 else DUSTY_ROSE
    add_oval(slide15, sx, sy, 0.75, 0.75, step_col)
    add_textbox(slide15, sx, sy, 0.75, 0.75, num, 22, "Georgia", CREAM,
                bold=True, align=PP_ALIGN.CENTER, margin_l=0, margin_t=0.16)
    add_textbox(slide15, sx, sy + 0.83, 3.0, 0.4, title, 15, "Georgia", BURGUNDY,
                bold=True, margin_t=0)
    add_textbox(slide15, sx, sy + 1.25, 3.0, 0.65, desc, 12, "Calibri", DARK_TEXT,
                margin_t=0)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Aplicaciones
# ══════════════════════════════════════════════════════════════════════════════
slide16 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide16, CREAM)
add_accent_bar(slide16)
add_header_band(slide16, "Aplicaciones Docentes e Investigadoras", DUSTY_ROSE)

add_rect(slide16, 0.3, 1.1, 4.4, 0.45, BURGUNDY)
add_textbox(slide16, 0.3, 1.1, 4.4, 0.45, "EN EL AULA", 14, "Georgia", CREAM,
            bold=True, align=PP_ALIGN.CENTER, margin_l=0, margin_t=0.1)
add_textbox_multiline(slide16, 0.3, 1.62, 4.4, 2.6, [
    ("Análisis comparativo de motetes y sus Misas parodia", False, False, True),
    ("Prácticas de análisis contrapuntístico asistido por ordenador", False, False, True),
    ("Proyectos de investigación sobre estilo y cita musical", False, False, True),
    ("Actividades de escucha activa con datos visuales", False, False, True),
], 13, "Calibri", DARK_TEXT)
add_oval(slide16, 0.5, 4.1, 1.6, 1.4, BURGUNDY, line_color=DUSTY_ROSE, line_width=1.5, transparency=80)

add_rect(slide16, 5.0, 1.1, 0.06, 4.1, DUSTY_ROSE)

add_rect(slide16, 5.28, 1.1, 4.4, 0.45, DUSTY_ROSE)
add_textbox(slide16, 5.28, 1.1, 4.4, 0.45, "EN LA INVESTIGACIÓN", 14, "Georgia", CREAM,
            bold=True, align=PP_ALIGN.CENTER, margin_l=0, margin_t=0.1)
add_textbox_multiline(slide16, 5.28, 1.62, 4.4, 2.6, [
    ("Identificación sistemática de citas y transformaciones", False, False, True),
    ("Estudios de recepción y transmisión musical", False, False, True),
    ("Análisis estadístico del estilo de un compositor", False, False, True),
    ("Publicaciones con visualizaciones reproducibles", False, False, True),
], 13, "Calibri", DARK_TEXT)
add_oval(slide16, 8.1, 4.1, 1.6, 1.4, DUSTY_ROSE, line_color=BURGUNDY, line_width=1.5, transparency=80)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Recursos
# ══════════════════════════════════════════════════════════════════════════════
slide17 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide17, CREAM)
add_accent_bar(slide17)
add_header_band(slide17, "Recursos y Comunidad")

resources17 = [
    (BURGUNDY,   "1", "crimproject.org",                             "Web principal del proyecto"),
    (DUSTY_ROSE, "2", "intervals-streamlit.crimproject.org",         "Aplicación web sin código"),
    (BURGUNDY,   "3", "github.com/HCDigitalScholarship/intervals",   "Código fuente y tutoriales"),
    (DUSTY_ROSE, "4", "pip install crim_intervals",                  "Biblioteca Python"),
    (BURGUNDY,   "5", "+50 investigadores",                          "Comunidad internacional activa"),
]
for i, (col, num, title, desc) in enumerate(resources17):
    ry = 1.15 + i * 0.7
    add_oval(slide17, 0.3, ry + 0.04, 0.52, 0.52, col)
    add_textbox(slide17, 0.3, ry + 0.04, 0.52, 0.52, num,
                14, "Georgia", CREAM, bold=True, align=PP_ALIGN.CENTER,
                margin_l=0, margin_t=0.1)
    add_textbox(slide17, 1.0, ry, 8.7, 0.6,
                title + "  —  " + desc,
                14, "Calibri", DARK_TEXT, margin_t=0.1)
    # Make title part bold by overwriting with rich text
    tb = slide17.shapes[-1]
    tf = tb.text_frame
    p = tf.paragraphs[0]
    # Clear and rebuild with bold/normal runs
    for run in p.runs:
        run.text = ""
    # Re-add
    r1 = p.add_run()
    r1.text = title
    r1.font.bold = True
    r1.font.size = pt_(14)
    r1.font.name = "Calibri"
    r1.font.color.rgb = DARK_TEXT
    r2 = p.add_run()
    r2.text = "  —  " + desc
    r2.font.bold = False
    r2.font.size = pt_(14)
    r2.font.name = "Calibri"
    r2.font.color.rgb = DARK_TEXT

add_rect(slide17, 0.3, 4.75, 9.4, 0.55, DARK_BURGUNDY)
add_textbox(slide17, 0.3, 4.75, 9.4, 0.55,
            "Los tutoriales en Jupyter Notebooks están disponibles gratuitamente en GitHub",
            13, "Calibri", CREAM, bold=True, italic=True, align=PP_ALIGN.CENTER,
            margin_l=0, margin_t=0.12)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — CONCLUSION (dark)
# ══════════════════════════════════════════════════════════════════════════════
slide18 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide18, DARK_BURGUNDY)
add_corner_ornament(slide18)
add_rect(slide18, 0, 0, 0.25, 5.625, DUSTY_ROSE)

add_textbox(slide18, 0.5, 0.28, 7.5, 0.65, "CRIM Intervals:",
            38, "Georgia", CREAM, bold=True)
add_textbox(slide18, 0.5, 0.93, 7.5, 0.55, "El Futuro del Análisis Musicológico",
            22, "Georgia", ROSE_LIGHT, italic=True)

bullets18 = [
    "Transforma el análisis de la polifonía renacentista",
    "Combina rigor analítico con accesibilidad (sin programación)",
    "Del análisis detallado al estudio masivo de corpus",
    "Herramienta ideal para estudiantes e investigadores",
]
for i, b in enumerate(bullets18):
    by = 1.65 + i * 0.6
    add_rect(slide18, 0.5, by + 0.12, 0.22, 0.22, DUSTY_ROSE)
    add_textbox(slide18, 0.9, by, 8.8, 0.5, b, 15, "Calibri", CREAM, margin_t=0.08)

add_rect(slide18, 0.5, 4.02, 9.0, 0.88, BURGUNDY)
shape18 = slide18.shapes[-1]
shape18.line.color.rgb = DUSTY_ROSE
shape18.line.width = Pt(1.5)
add_textbox(slide18, 0.5, 4.02, 9.0, 0.88,
            '"Las Humanidades Digitales al servicio del repertorio histórico"',
            17, "Georgia", CREAM, italic=True, align=PP_ALIGN.CENTER,
            margin_l=0, margin_t=0.22)

add_textbox(slide18, 0.5, 5.2, 9.0, 0.3,
            "intervals-streamlit.crimproject.org  ·  crimproject.org",
            11, "Calibri", ROSE_LIGHT, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUTPUT_PATH)
print(f"Saved: {OUTPUT_PATH}")
